"""
PowerPoint to Markdown converter with text and image extraction.
"""

import base64
from datetime import datetime
from pathlib import Path
from typing import Optional
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

from .image_optimizer import ImageOptimizer


class PowerPointConverter:
    """Convert PowerPoint files (.pptx, .ppt) to Markdown format."""

    def __init__(self, cache_dir: Path, image_optimizer: ImageOptimizer):
        """
        Initialize the PowerPoint converter.

        Args:
            cache_dir: Directory for caching outputs
            image_optimizer: Image optimization utility
        """
        self.cache_dir = cache_dir
        self.image_optimizer = image_optimizer

    def convert_file(
        self,
        file_path: str,
        extract_images: bool = True,
        image_format: str = "file",
        output_name: Optional[str] = None,
        conversion_dir: Path = None,
        images_dir: Path = None,
    ) -> dict:
        """
        Convert PowerPoint file to Markdown.

        Args:
            file_path: Path to PowerPoint file
            extract_images: Whether to extract embedded images
            image_format: "file", "base64", or "both"
            output_name: Custom output name
            conversion_dir: Directory for this conversion
            images_dir: Directory for extracted images

        Returns:
            dict with markdown content, images, and metadata
        """
        file_path = Path(file_path)

        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        suffix = file_path.suffix.lower()
        if suffix not in [".pptx", ".ppt"]:
            raise ValueError(f"Unsupported file format: {suffix}")

        # Load presentation
        prs = Presentation(file_path)

        markdown_parts = []
        all_images = []
        image_counter = 0

        # Add document title
        markdown_parts.append(f"# PowerPoint: {file_path.name}\n")
        markdown_parts.append(f"**Total slides: {len(prs.slides)}**\n")
        markdown_parts.append("---\n")

        # Process each slide
        for slide_num, slide in enumerate(prs.slides, 1):
            markdown_parts.append(f"\n## Slide {slide_num}\n")

            # Get slide layout name if available
            if slide.slide_layout:
                layout_name = slide.slide_layout.name
                markdown_parts.append(f"*Layout: {layout_name}*\n")

            # Process shapes on the slide
            slide_text_parts = []
            slide_tables = []

            for shape in slide.shapes:
                # Handle text frames
                if shape.has_text_frame:
                    text = self._extract_text_from_shape(shape)
                    if text.strip():
                        slide_text_parts.append(text)

                # Handle tables
                if shape.has_table:
                    table_md = self._extract_table(shape.table)
                    if table_md:
                        slide_tables.append(table_md)

                # Handle images
                if extract_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        image_counter += 1
                        img_name = f"slide{slide_num}_image_{image_counter:03d}"

                        # Get image data
                        image = shape.image
                        img_bytes = image.blob
                        pil_img = Image.open(BytesIO(img_bytes))

                        # Optimize image
                        optimized_data, ext = self.image_optimizer.optimize_image(
                            pil_img, img_name
                        )

                        img_filename = f"{img_name}{ext}"
                        img_path = images_dir / img_filename

                        if image_format in ["file", "both"]:
                            with open(img_path, "wb") as f:
                                f.write(optimized_data)
                            all_images.append(str(img_path))
                            slide_text_parts.append(f"\n![{img_name}](images/{img_filename})")

                        if image_format in ["base64", "both"]:
                            b64_data = base64.b64encode(optimized_data).decode("utf-8")
                            content_type = f"image/{ext.lstrip('.')}"
                            if image_format == "base64":
                                slide_text_parts.append(
                                    f"\n![{img_name}](data:{content_type};base64,{b64_data})"
                                )
                    except Exception as e:
                        # Log but continue on image extraction errors
                        slide_text_parts.append(f"\n*[Image extraction failed: {e}]*")

                # Handle grouped shapes
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    group_text = self._extract_from_group(shape)
                    if group_text.strip():
                        slide_text_parts.append(group_text)

            # Add extracted content
            if slide_text_parts:
                markdown_parts.append("\n".join(slide_text_parts))

            if slide_tables:
                markdown_parts.append("\n### Tables\n")
                for table in slide_tables:
                    markdown_parts.append(table)

            # Add speaker notes if present
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                if notes_frame and notes_frame.text.strip():
                    markdown_parts.append("\n**Speaker Notes:**")
                    markdown_parts.append(f"> {notes_frame.text.strip()}")

            markdown_parts.append("\n---\n")

        # Extract metadata
        metadata = self._extract_metadata(prs, file_path)

        markdown = "\n".join(markdown_parts)

        return {
            "markdown": markdown,
            "images": all_images,
            "metadata": metadata,
        }

    def _extract_text_from_shape(self, shape) -> str:
        """Extract formatted text from a shape."""
        if not shape.has_text_frame:
            return ""

        text_parts = []

        for paragraph in shape.text_frame.paragraphs:
            para_text = ""

            for run in paragraph.runs:
                text = run.text
                if not text:
                    continue

                # Apply formatting
                if run.font.bold:
                    text = f"**{text}**"
                if run.font.italic:
                    text = f"*{text}*"
                if run.font.underline:
                    text = f"<u>{text}</u>"

                para_text += text

            if para_text.strip():
                # Check paragraph level for list indentation
                level = paragraph.level or 0
                if level > 0:
                    indent = "  " * level
                    para_text = f"{indent}- {para_text}"

                text_parts.append(para_text)

        return "\n".join(text_parts)

    def _extract_from_group(self, group_shape) -> str:
        """Recursively extract text from grouped shapes."""
        text_parts = []

        for shape in group_shape.shapes:
            if shape.has_text_frame:
                text = self._extract_text_from_shape(shape)
                if text.strip():
                    text_parts.append(text)
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                nested_text = self._extract_from_group(shape)
                if nested_text.strip():
                    text_parts.append(nested_text)

        return "\n".join(text_parts)

    def _extract_table(self, table) -> str:
        """Convert PowerPoint table to Markdown format."""
        if not table.rows:
            return ""

        md_rows = []

        for i, row in enumerate(table.rows):
            cells = []
            for cell in row.cells:
                # Get cell text
                cell_text = cell.text.strip().replace("\n", " ").replace("|", "\\|")
                cells.append(cell_text)

            md_rows.append("| " + " | ".join(cells) + " |")

            # Add header separator after first row
            if i == 0:
                separator = "| " + " | ".join(["---"] * len(cells)) + " |"
                md_rows.append(separator)

        return "\n".join(md_rows) + "\n"

    def _extract_metadata(self, presentation, file_path: Path) -> dict:
        """Extract PowerPoint metadata."""
        props = presentation.core_properties

        return {
            "source": str(file_path),
            "converted_at": datetime.now().isoformat(),
            "title": props.title or "",
            "author": props.author or "",
            "created": str(props.created) if props.created else "",
            "modified": str(props.modified) if props.modified else "",
            "subject": props.subject or "",
            "slide_count": len(presentation.slides),
            "slide_width": presentation.slide_width.inches if presentation.slide_width else None,
            "slide_height": presentation.slide_height.inches if presentation.slide_height else None,
        }
